from pathlib import Path
import re

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent / "SoftwareTestingDemo" / "SoftwareTestingDemo"
OUTPUT_DIR = BASE_DIR / "output"
ASSETS_DIR = OUTPUT_DIR / "assets"
OUTPUT_PPT = OUTPUT_DIR / "PJ1测试汇报_成果展示风格版.pptx"
FALLBACK_OUTPUT_PPT = OUTPUT_DIR / "PJ1测试汇报_成果展示风格版_优化版.pptx"

FONT = "Microsoft YaHei"
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 247, 250)
GRAY = RGBColor(105, 105, 105)
GREEN = RGBColor(114, 178, 46)
DARK_GREEN = RGBColor(16, 157, 113)
CYAN = RGBColor(41, 173, 209)
BLUE = RGBColor(42, 128, 185)
DEEP_BLUE = RGBColor(31, 97, 141)
ORANGE = RGBColor(241, 145, 0)
BROWN = RGBColor(171, 72, 36)

STRIP_COLORS = [
    RGBColor(0, 156, 104),
    RGBColor(122, 192, 57),
    RGBColor(62, 174, 203),
    RGBColor(0, 137, 195),
    RGBColor(48, 107, 163),
    RGBColor(213, 82, 31),
    RGBColor(247, 148, 29),
]


def ensure_dirs() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def setup_chart_font() -> None:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False


def create_coverage_chart(path: Path) -> None:
    setup_chart_font()
    labels = ["指令覆盖率", "行覆盖率", "分支覆盖率", "方法覆盖率", "类覆盖率"]
    values = [92.54, 92.34, 80.77, 93.53, 95.24]
    colors = ["#009C68", "#7AC039", "#3EAECB", "#0089C3", "#F7941D"]

    fig, ax = plt.subplots(figsize=(8.0, 3.4), dpi=180)
    bars = ax.barh(labels, values, color=colors, height=0.55)
    ax.set_xlim(0, 100)
    ax.invert_yaxis()
    ax.grid(axis="x", linestyle="--", alpha=0.22)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#BDBDBD")
    ax.tick_params(axis="y", labelsize=10)
    ax.tick_params(axis="x", labelsize=9)
    for bar, value in zip(bars, values):
        ax.text(value + 1.2, bar.get_y() + bar.get_height() / 2, f"{value:.2f}%", va="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = WHITE, width: float = 0.5) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def style_text_frame(shape, size=20, bold=False, color=BLACK, align=None, auto_size=False) -> None:
    tf = shape.text_frame
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in tf.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_title(slide, text: str, left=0.82, top=0.35, width=11.5, height=0.7, size=30) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text = text
    style_text_frame(box, size=size, bold=True)


def add_footer_strip(slide) -> None:
    start_left = Inches(10.62)
    top = Inches(7.05)
    block_w = Inches(0.16)
    gap = Inches(0.035)
    for idx, color in enumerate(STRIP_COLORS):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            start_left + idx * (block_w + gap),
            top,
            block_w,
            Inches(0.09),
        )
        set_fill(rect, color)
        rect.line.fill.background()


def add_small_caption(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.25), Inches(7.02), Inches(1.4), Inches(0.18))
    box.text = text
    style_text_frame(box, size=7, color=BLACK)


def add_textbox(slide, text: str, left, top, width, height, size=18, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text = text
    style_text_frame(box, size=size, bold=bold, color=color, align=align, auto_size=True)
    return box


def add_side_title(slide, text: str, left=0.25, top=0.58, width=0.55, height=5.15, size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text = "\n".join(text)
    tf = box.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 0.85
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = BLACK
    return box


def add_bullets(slide, lines, left, top, width, height, size=17, heading=None):
    text = ""
    if heading:
        text += f"{heading}\n"
    text += "\n".join(f"• {line}" for line in lines)
    box = add_textbox(slide, text, left, top, width, height, size=size)
    if heading:
        first = box.text_frame.paragraphs[0]
        for run in first.runs:
            run.font.bold = True
            run.font.size = Pt(size + 2)
    return box


def add_table(slide, rows, left, top, width, height, header_color=GREEN, alt_color=RGBColor(229, 241, 216), font_size=12):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            fill_color = header_color if row_idx == 0 else (alt_color if row_idx % 2 else RGBColor(241, 247, 236))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if row_idx else font_size + 1)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = WHITE if row_idx == 0 else BLACK
    return table_shape


def add_metric_card(slide, value: str, label: str, left, top, color, width=2.2):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.0))
    set_fill(card, color)
    set_line(card, color)
    card.text = f"{value}\n{label}"
    style_text_frame(card, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_module_card(slide, title: str, body: str, left, top, width, height, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    set_fill(rect, RGBColor(246, 248, 249))
    set_line(rect, RGBColor(230, 230, 230))
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42))
    set_fill(head, color)
    head.line.fill.background()
    head.text = title
    style_text_frame(head, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, body, left + 0.08, top + 0.52, width - 0.16, height - 0.62, size=10, color=BLACK, align=PP_ALIGN.CENTER)


def add_code_block(slide, text: str, left, top, width, height):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    set_fill(box, RGBColor(247, 249, 252))
    set_line(box, RGBColor(213, 220, 228))
    box.text = text
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(40, 60, 80)


def extract_java_method_snippet(file_path: Path, method_name: str) -> str:
    content = file_path.read_text(encoding="utf-8")
    pattern = re.compile(
        rf"(^\s*@Test\s*\n)?^\s*void\s+{re.escape(method_name)}\s*\([^)]*\)\s*(?:throws [^{{]+)?\{{",
        re.MULTILINE,
    )
    match = pattern.search(content)
    if not match:
        raise ValueError(f"Method not found: {method_name} in {file_path}")

    start = match.start()
    brace_start = content.find("{", match.end() - 1)
    depth = 0
    end = brace_start
    for idx in range(brace_start, len(content)):
        char = content[idx]
        if char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                end = idx
                break

    return content[start : end + 1].strip()


def new_slide(prs: Presentation, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_footer_strip(slide)
    add_small_caption(slide, f"图表{number}")
    return slide


def build_deck() -> None:
    ensure_dirs()
    coverage_chart = ASSETS_DIR / "coverage_summary_showcase.png"
    create_coverage_chart(coverage_chart)
    integration_snippet = extract_java_method_snippet(
        PROJECT_DIR / "src" / "test" / "java" / "com" / "demo" / "controller" / "OrderControllerIntegrationTest.java",
        "shouldReturnLoggedUserOrdersOnly",
    )
    unit_snippet = extract_java_method_snippet(
        PROJECT_DIR / "src" / "test" / "java" / "com" / "demo" / "service" / "impl" / "UserServiceImplTest.java",
        "shouldReturnNewUserCountAfterCreate",
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = new_slide(prs, 1)
    add_textbox(slide, "SoftwareTestingDemo  测试成果展示", 1.45, 1.2, 10.5, 0.8, size=29, bold=True, align=PP_ALIGN.CENTER)
    subtitle = "第 PJ1 小组   胡彬泽 沙华煜 彭祈元 郭政颍 周宇尘\n测试计划标识   Project1_TestPlan_V1.2_20260423"
    add_textbox(slide, subtitle, 1.2, 3.4, 10.9, 0.9, size=17, align=PP_ALIGN.CENTER)

    # 2
    slide = new_slide(prs, 2)
    add_title(slide, "PJ1主要内容")
    add_bullets(
        slide,
        ["测试项与先决条件", "通过/失败准则与风险管控", "小组分工与时间安排", "交付物与可复现环境"],
        0.9,
        1.8,
        5.5,
        3.9,
        heading="1. 测试计划规定",
    )
    add_bullets(
        slide,
        ["service 单元测试与 controller 集成测试", "黑盒等价类/边界值 + 白盒语句/判定覆盖", "JUnit5、Mockito、MockMvc、H2、JaCoCo", "执行结果、覆盖率与测试小结"],
        6.7,
        1.8,
        5.5,
        3.9,
        heading="2. 测试本身完成",
    )

    # 3
    slide = new_slide(prs, 3)
    add_title(slide, "测试项 & 先决条件")
    add_bullets(
        slide,
        ["service 层：7 个实现类，覆盖核心业务方法 45 个", "controller 层：11 类控制器，覆盖 65 个请求映射", "测试脚本：7 个 service 测试类 + 8 个 controller 集成测试类", "测试入口：项目根目录执行 mvn test"],
        0.8,
        1.6,
        5.8,
        3.2,
        size=15,
        heading="测试项版本/修订号",
    )
    add_bullets(
        slide,
        ["JDK 8+、Maven、Spring Boot 2.2.2", "H2 内存库启用 MySQL 兼容模式，不依赖本机 MySQL", "JUnit5 + Mockito + MockMvc 作为自动化测试基座", "JaCoCo 报告生成到 target/site/jacoco/index.html"],
        7.0,
        1.6,
        5.2,
        3.2,
        size=15,
        heading="其中的环境需求",
    )

    # 4
    slide = new_slide(prs, 4)
    add_bullets(
        slide,
        ["全部自动化用例通过", "JaCoCo 报告正常生成", "整体指令覆盖率 ≥ 85%，分支覆盖率 ≥ 75%", "新克隆环境可重复执行"],
        0.9,
        1.05,
        5.3,
        4.9,
        size=15,
        heading="测试项通过/失败准则",
    )
    add_bullets(
        slide,
        ["测试基座或依赖失效时暂停", "核心模块无法编译时暂停", "严重缺陷修复、环境恢复后继续", "需求或范围变化后同步更新用例与文档"],
        6.8,
        1.05,
        5.3,
        4.9,
        size=15,
        heading="暂停准则和继续准则",
    )
    add_textbox(slide, "通过：128 条测试全部通过；失败：出现阻断缺陷、报告缺失或环境不可复现。", 0.95, 6.0, 10.8, 0.45, size=14, color=BROWN)

    # 5
    slide = new_slide(prs, 5)
    add_title(slide, "风险管控")
    add_table(
        slide,
        [
            ["风险类型", "本项目表现", "应急措施"],
            ["数据库差异", "H2 与 MySQL 可能存在少量行为差异", "启用 MySQL 模式；差异写入测试总结"],
            ["覆盖盲区", "FileUtil 真实落盘未完全覆盖", "列为已知限制；后续隔离文件系统"],
            ["脚本维护", "文档和脚本可能不同步", "统一编号；提交前执行 mvn test"],
            ["UI 风险", "未做浏览器级 UI 自动化", "明确范围外；后续补 E2E 回归"],
        ],
        1.55,
        1.7,
        10.2,
        3.4,
        header_color=CYAN,
        alt_color=RGBColor(222, 239, 246),
        font_size=12,
    )

    # 6
    slide = new_slide(prs, 6)
    add_textbox(slide, "组员\n分工", 0.35, 0.5, 1.4, 1.4, size=27, bold=True)
    add_table(
        slide,
        [
            ["序号", "小组成员", "职责", "任务占比"],
            ["1", "胡彬泽", "环境搭建；service 单测；controller 集成测试；覆盖率报告", "40%"],
            ["2", "沙华煜", "项目测试计划书；统一团队口径；交付清单", "15%"],
            ["3", "彭祈元", "单元测试用例文档；维护 UT 编号", "15%"],
            ["4", "郭政颍", "集成测试用例文档；维护 controller 编号", "15%"],
            ["5", "周宇尘", "测试总结；结果汇总；改进建议", "15%"],
        ],
        2.35,
        0.65,
        9.5,
        5.4,
        header_color=DARK_GREEN,
        alt_color=RGBColor(218, 236, 226),
        font_size=11,
    )

    # 7
    slide = new_slide(prs, 7)
    add_title(slide, "时间安排", left=0.25, top=0.18)
    add_table(
        slide,
        [
            ["日期", "任务安排", "输出成果"],
            ["04.06-04.12", "service 单元测试编写", "7 个 service 测试类"],
            ["04.10-04.16", "controller 集成测试编写", "8 个集成测试脚本"],
            ["04.17-04.18", "全量 mvn test 与覆盖率收集", "JaCoCo + Surefire 报告"],
            ["04.19-04.22", "测试计划、用例文档、总结对齐", "三份课程文档"],
            ["04.23-04.26", "PPT 汇报与提交包整理", "PDF + 源码压缩包"],
        ],
        0.45,
        1.25,
        12.0,
        4.75,
        header_color=GREEN,
        alt_color=RGBColor(226, 240, 210),
        font_size=13,
    )

    # 8
    slide = new_slide(prs, 8)
    add_textbox(slide, "PJ1测试的主要方法", 0.55, 0.6, 5.3, 0.45, size=19, bold=True)
    add_bullets(
        slide,
        ["主要活动：编写测试脚本、执行测试用例、分析覆盖率与结果", "技术：JUnit5、SpringBootTest、Mockito、MockMvc、H2、JaCoCo", "黑盒：等价类划分、边界值分析，用于参数、ID、状态值、页码等输入", "白盒：语句覆盖、判定覆盖，用于 service 分支与 controller 登录/跳转逻辑", "冗余控制：合并同等价类和相同结果路径，保留代表性用例"],
        0.7,
        1.35,
        11.8,
        4.7,
        size=15,
    )

    # 9
    slide = new_slide(prs, 9)
    add_title(slide, "测试覆盖 service 与 controller 两层核心模块")
    modules = [
        ("用户与登录", "UserService\nUserController\nAdminUserController", DARK_GREEN),
        ("场馆信息", "VenueService\nVenueController\nAdminVenueController", GREEN),
        ("预约订单", "OrderService\nOrderVoService\nOrderController", CYAN),
        ("留言板", "MessageService\nMessageVoService\nMessageController", BLUE),
        ("新闻展示", "NewsService\nNewsController\nAdminNewsController", ORANGE),
        ("公共首页", "IndexController\n公共内容展示\n首页数据聚合", BROWN),
    ]
    for idx, (title, body, color) in enumerate(modules):
        add_module_card(slide, title, body, 0.65 + idx * 2.05, 1.9, 1.76, 3.1, color)

    # 10
    slide = new_slide(prs, 10)
    add_title(slide, "不会被测试的特性")
    add_table(
        slide,
        [
            ["序号", "不会被测试的特性", "原因"],
            ["1", "浏览器 UI 自动化", "PJ1 要求聚焦 service 单元测试与 controller 集成测试"],
            ["2", "性能压测", "课程任务未要求吞吐、并发与响应时间指标"],
            ["3", "安全渗透测试", "仅验证登录态与权限分支，不做漏洞扫描"],
            ["4", "真实 MySQL 专有行为", "H2 MySQL 模式用于可复现测试，真实库差异列为限制"],
            ["5", "真实文件落盘", "FileUtil 与业务主流程关联较弱，后续可单独隔离补测"],
        ],
        0.55,
        1.55,
        12.0,
        3.9,
        header_color=GREEN,
        alt_color=RGBColor(226, 240, 210),
        font_size=12,
    )

    # 11
    slide = new_slide(prs, 11)
    add_title(slide, "测试用例设计文档\n以订单状态流转为例", top=0.2, size=25)
    add_table(
        slide,
        [
            ["设计技术", "代表输入/场景", "预期结果", "编号口径"],
            ["等价类划分", "合法用户、合法场馆、存在的订单 ID", "订单创建/查询成功", "订单相关用例"],
            ["边界值分析", "page 默认值、非法页码、空参数", "返回默认页或异常路径被验证", "分页与参数用例"],
            ["语句覆盖", "创建、取消、查询等核心业务路径", "可达语句至少执行一次", "单元测试用例"],
            ["判定覆盖", "审核通过/拒绝、未登录/已登录", "true/false 两侧均被验证", "集成测试用例"],
        ],
        0.45,
        1.65,
        12.15,
        4.25,
        header_color=ORANGE,
        alt_color=RGBColor(252, 225, 183),
        font_size=11,
    )

    # 12
    slide = new_slide(prs, 12)
    add_title(slide, "覆盖率")
    add_metric_card(slide, "128", "自动化测试全部通过", 0.55, 1.45, DARK_GREEN)
    add_metric_card(slide, "92.54%", "整体指令覆盖率", 0.55, 2.75, CYAN)
    add_metric_card(slide, "80.77%", "整体分支覆盖率", 0.55, 4.05, ORANGE)
    slide.shapes.add_picture(str(coverage_chart), Inches(3.5), Inches(1.25), Inches(8.6), Inches(4.25))
    add_textbox(slide, "补充：service.impl 指令/分支覆盖率 100%；controller 指令覆盖率 98.54%，分支覆盖率 82.35%。", 3.7, 5.7, 8.2, 0.45, size=13, color=GRAY)

    # 13
    slide = new_slide(prs, 13)
    add_title(slide, "每个测试类里的基本思路")
    add_bullets(
        slide,
        ["service 单测：@ExtendWith + MockitoExtension 隔离 DAO", "controller 集成：@SpringBootTest + @AutoConfigureMockMvc 启动上下文", "BeforeEach / TestDataFactory 统一造数，保证测试互不污染", "MockMvc perform 后校验 status、view、redirect、model 属性", "执行 mvn test 后由 Surefire 与 JaCoCo 汇总结果"],
        0.7,
        1.35,
        11.9,
        4.8,
        size=15,
    )

    # 14
    slide = new_slide(prs, 14)
    add_title(slide, "示例")
    add_code_block(
        slide,
        integration_snippet,
        0.7,
        1.35,
        5.8,
        3.55,
    )
    add_code_block(
        slide,
        unit_snippet,
        6.8,
        1.35,
        5.8,
        3.55,
    )
    add_textbox(slide, "左侧：真实集成测试脚本片段\n右侧：真实单元测试脚本片段", 2.6, 5.35, 8.2, 0.55, size=16, color=BROWN, align=PP_ALIGN.CENTER)

    # 15
    slide = new_slide(prs, 15)
    add_side_title(slide, "单元测试模块", size=22)
    add_table(
        slide,
        [
            ["序号", "测试模块", "用例数量", "覆盖重点"],
            ["1", "用户服务", "8", "登录、注册、增删改查"],
            ["2", "场馆服务", "8", "场馆查询、更新、删除"],
            ["3", "订单服务", "14", "订单创建、取消、审核状态"],
            ["4", "留言服务", "11", "留言发布、审核、删除"],
            ["5", "新闻服务", "5", "新闻列表、详情、维护"],
            ["6", "聚合查询服务", "4", "VO 聚合查询与边界路径"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=DARK_GREEN,
        alt_color=RGBColor(218, 236, 226),
        font_size=12,
    )

    # 16
    slide = new_slide(prs, 16)
    add_side_title(slide, "公共展示模块", size=22)
    add_table(
        slide,
        [
            ["序号", "特性", "代表场景", "测试设计技术"],
            ["1", "首页展示", "首页内容正常加载", "集成测试"],
            ["2", "新闻列表", "按页展示新闻数据", "等价类 + 边界值"],
            ["3", "新闻详情", "查看存在 / 不存在新闻", "边界值"],
            ["4", "场馆列表", "展示场馆与分页信息", "等价类"],
            ["5", "预约入口", "登录与未登录访问分流", "判定覆盖"],
            ["6", "场馆查询", "按名称查询有效/无效场馆", "等价类"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=GREEN,
        alt_color=RGBColor(226, 240, 210),
        font_size=12,
    )

    # 17
    slide = new_slide(prs, 17)
    add_side_title(slide, "预约订单模块", size=22)
    add_table(
        slide,
        [
            ["序号", "特性", "代表场景", "覆盖重点"],
            ["1", "订单列表", "分页查看个人订单", "登录态 + page 边界"],
            ["2", "创建订单", "合法预约 / 无效输入", "有效类 / 无效类"],
            ["3", "取消订单", "本人订单与非本人订单", "权限与状态流转"],
            ["4", "管理员列表", "按状态筛选订单", "筛选与分页"],
            ["5", "审核订单", "通过 / 拒绝两类结果", "判定覆盖 true/false"],
            ["6", "状态推进", "创建、取消、审核全流程", "语句覆盖 + 异常路径"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=CYAN,
        alt_color=RGBColor(219, 240, 247),
        font_size=12,
    )

    # 18
    slide = new_slide(prs, 18)
    add_side_title(slide, "用户登录模块", size=22)
    add_table(
        slide,
        [
            ["序号", "特性", "代表场景", "测试设计技术"],
            ["1", "登录成功", "合法账号与密码", "有效等价类"],
            ["2", "登录失败", "错误密码 / 不存在账号", "无效等价类"],
            ["3", "注册用户", "新账号 / 重复账号", "边界值 + 重复值"],
            ["4", "退出登录", "清理会话状态", "状态校验"],
            ["5", "管理员维护", "用户增删改查", "判定覆盖"],
            ["6", "服务校验", "登录、注册、删除、更新", "语句覆盖"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=BLUE,
        alt_color=RGBColor(218, 232, 245),
        font_size=12,
    )

    # 19
    slide = new_slide(prs, 19)
    add_side_title(slide, "留言板模块", size=23)
    add_table(
        slide,
        [
            ["序号", "特性", "代表场景", "覆盖重点"],
            ["1", "留言列表", "分页查看留言", "分页边界"],
            ["2", "发布留言", "登录用户发布有效内容", "登录态 + 有效输入"],
            ["3", "空内容留言", "空内容 / 缺少参数", "无效等价类"],
            ["4", "删除留言", "本人删除 / 非法删除", "权限路径"],
            ["5", "管理员审核", "通过 / 拒绝留言", "判定覆盖"],
            ["6", "聚合查询", "留言与用户信息组合", "聚合查询"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=BROWN,
        alt_color=RGBColor(246, 222, 212),
        font_size=12,
    )

    # 20
    slide = new_slide(prs, 20)
    add_side_title(slide, "后台管理模块", size=22)
    add_table(
        slide,
        [
            ["序号", "特性", "代表场景", "覆盖重点"],
            ["1", "场馆新增/修改", "合法数据 / 边界价格", "有效类 + 边界值"],
            ["2", "场馆删除", "存在 ID / 不存在 ID", "异常路径"],
            ["3", "新闻新增/修改", "空标题 / 有效内容", "无效类 + 有效类"],
            ["4", "新闻删除", "删除后重定向", "状态码与跳转"],
            ["5", "用户管理", "重复账号与无效 ID", "重复值 + 边界值"],
            ["6", "公共跳转", "视图、模型、重定向", "集成链路验证"],
        ],
        1.05,
        0.98,
        11.65,
        5.65,
        header_color=ORANGE,
        alt_color=RGBColor(252, 225, 183),
        font_size=12,
    )

    # 21
    slide = new_slide(prs, 21)
    add_textbox(slide, "测试小结", 0.65, 0.55, 4.8, 0.75, size=28, bold=True)
    add_bullets(
        slide,
        [
            "本次测试已覆盖 service 与 controller 的核心功能链路，自动化脚本可重复执行。",
            "128 条测试全部通过，整体指令覆盖率 92.54%，分支覆盖率 80.77%，满足测试计划出口准则。",
            "用例设计同时体现等价类划分、边界值分析、语句覆盖与判定覆盖，兼顾完整性与冗余控制。",
            "当前范围不包含浏览器 UI 自动化、性能测试和真实 MySQL 差异验证，这些已在《测试总结》中列为后续优化方向。",
            "提交前建议再次执行 mvn test，并将三份文档 PDF 与源码统一打包为 PJ1-小组编号.zip。",
        ],
        1.0,
        1.75,
        11.0,
        4.6,
        size=15,
    )

    try:
        prs.save(OUTPUT_PPT)
    except PermissionError:
        prs.save(FALLBACK_OUTPUT_PPT)


if __name__ == "__main__":
    build_deck()
